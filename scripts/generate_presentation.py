from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_FILE = "ThuyetTrinh_QuanLyPhongKhamBenh_ChiTiet.pptx"

# Color palette
BG_LIGHT = RGBColor(245, 248, 252)
BG_DARK = RGBColor(20, 45, 78)
ACCENT = RGBColor(10, 132, 255)
ACCENT_2 = RGBColor(0, 166, 153)
TEXT_DARK = RGBColor(24, 24, 24)
TEXT_LIGHT = RGBColor(250, 250, 250)
MUTED = RGBColor(96, 96, 96)


def set_background(slide, dark=False):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK if dark else BG_LIGHT


def add_top_bar(slide, title, subtitle=None, dark=False):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT if not dark else ACCENT_2
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(9.8), Inches(0.5))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Calibri"
    p.font.bold = True
    p.font.size = Pt(25)
    p.font.color.rgb = TEXT_LIGHT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.45), Inches(1.0), Inches(12.3), Inches(0.5))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Calibri"
        sp.font.size = Pt(15)
        sp.font.color.rgb = MUTED if not dark else RGBColor(220, 220, 220)


def add_bullets(slide, items, left=0.7, top=1.7, width=12.0, height=5.2, dark=False, level0_size=22):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, item in enumerate(items):
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.name = "Calibri"
        p.font.size = Pt(level0_size - min(level * 3, 8))
        p.space_after = Pt(8)
        p.font.color.rgb = TEXT_LIGHT if dark else TEXT_DARK


def add_two_columns(slide, left_title, left_items, right_title, right_items, dark=False):
    left_header = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(5.9), Inches(0.4))
    ltf = left_header.text_frame
    ltf.clear()
    lp = ltf.paragraphs[0]
    lp.text = left_title
    lp.font.bold = True
    lp.font.size = Pt(18)
    lp.font.name = "Calibri"
    lp.font.color.rgb = ACCENT if not dark else ACCENT_2

    right_header = slide.shapes.add_textbox(Inches(6.85), Inches(1.45), Inches(5.9), Inches(0.4))
    rtf = right_header.text_frame
    rtf.clear()
    rp = rtf.paragraphs[0]
    rp.text = right_title
    rp.font.bold = True
    rp.font.size = Pt(18)
    rp.font.name = "Calibri"
    rp.font.color.rgb = ACCENT if not dark else ACCENT_2

    add_bullets(slide, left_items, left=0.7, top=1.9, width=5.9, height=4.9, dark=dark, level0_size=18)
    add_bullets(slide, right_items, left=6.85, top=1.9, width=5.9, height=4.9, dark=dark, level0_size=18)


def add_footer(slide, text):
    footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.08), Inches(12.4), Inches(0.28))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(130, 130, 130)
    p.alignment = PP_ALIGN.RIGHT


def add_process_line(slide, steps):
    top = Inches(2.6)
    left_start = 0.85
    step_width = 2.35
    gap = 0.4

    for idx, step in enumerate(steps):
        x = Inches(left_start + idx * (step_width + gap))
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, Inches(step_width), Inches(1.1))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(234, 243, 255) if idx % 2 == 0 else RGBColor(224, 245, 241)
        rect.line.color.rgb = ACCENT

        t = rect.text_frame
        t.clear()
        p = t.paragraphs[0]
        p.text = step
        p.font.name = "Calibri"
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER

        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + Inches(step_width + 0.06),
                top + Inches(0.37),
                Inches(0.28),
                Inches(0.36),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()


def add_slide(prs, title, subtitle=None, bullets=None, dark=False, footer="Đồ án Quản Lý Phòng Khám Bệnh - Nhóm thực hiện"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, dark=dark)
    add_top_bar(slide, title, subtitle, dark=dark)
    if bullets:
        add_bullets(slide, bullets, dark=dark)
    add_footer(slide, footer)
    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1 - Title
    slide = add_slide(
        prs,
        "QUẢN LÝ PHÒNG KHÁM BỆNH",
        "Báo cáo phân tích hệ thống, kiến trúc kỹ thuật và kịch bản demo thực tế",
        dark=True,
    )
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.35), Inches(5.8), Inches(1.05))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(34, 74, 124)
    badge.line.color.rgb = ACCENT_2
    tf = badge.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Desktop App Java Swing - RBAC động - Nghiệp vụ phòng khám + nhà thuốc"
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # 2 - Agenda
    add_slide(
        prs,
        "Nội dung thuyết trình",
        bullets=[
            "Bối cảnh bài toán và mục tiêu sản phẩm",
            "Kiến trúc hệ thống theo lớp GUI-BUS-DAO-DB",
            "Cơ chế phân quyền động RBAC và bảo mật truy cập",
            "4 luồng nghiệp vụ trọng tâm theo vai trò người dùng",
            "Thiết kế dữ liệu và quan hệ bảng chính",
            "Điểm mạnh kỹ thuật, hạn chế và hướng phát triển",
            "Kịch bản demo thuyết phục + bộ câu hỏi phản biện",
        ],
    )

    # 3 - Problem/Motivation
    add_slide(
        prs,
        "Bối cảnh và mục tiêu",
        bullets=[
            "Phòng khám cần một hệ thống tập trung để quản lý khám bệnh, lịch hẹn, bán thuốc và kho.",
            "Mỗi nhóm người dùng có quyền hạn khác nhau: Khách, Bác sĩ, Nhà thuốc, Quản trị.",
            "Mục tiêu chính: số hóa quy trình đầu-cuối, giảm sai sót thủ công, kiểm soát trạng thái nghiệp vụ.",
            "Mục tiêu kỹ thuật: phân quyền linh hoạt theo dữ liệu, không phải sửa code khi đổi quyền.",
        ],
    )

    # 4 - Product scope
    add_slide(
        prs,
        "Phạm vi hệ thống",
        bullets=[
            "Guest: đặt lịch khám, mua thuốc lẻ, in chứng từ.",
            "Doctor: đăng ký lịch làm, xử lý lịch khám, lập bệnh án, kê đơn, xác nhận hóa đơn khám.",
            "Pharmacist: quản lý thuốc/NCC, nhập kho theo lô-hạn dùng, bán thuốc và truy vết xuất kho.",
            "Admin: dashboard, quản lý danh mục, tài khoản, phân quyền role-permission theo thời gian thực.",
        ],
    )

    # 5 - Architecture
    slide = add_slide(prs, "Kiến trúc tổng thể", "Mô hình nhiều lớp rõ trách nhiệm")
    add_process_line(slide, ["GUI", "BUS", "DAO", "DB"])
    add_bullets(
        slide,
        [
            "GUI: hiển thị panel theo route, chặn truy cập khi thiếu quyền.",
            "BUS: kiểm tra nghiệp vụ, chuẩn hóa trạng thái, điều phối nhiều DAO.",
            "DAO: thao tác SQL có tham số PreparedStatement, ánh xạ ResultSet sang DTO.",
            "DB: MySQL + HikariCP để tối ưu kết nối và ổn định khi tải tăng.",
        ],
        top=4.2,
        height=2.6,
        level0_size=17,
    )

    # 6 - Package structure
    slide = add_slide(prs, "Cấu trúc mã nguồn", "Tổ chức package bám đúng domain")
    add_two_columns(
        slide,
        "Gói lõi",
        [
            "phongkham.gui: giao diện theo module admin/bacsi/guest/nhathuoc",
            "phongkham.BUS: xử lý nghiệp vụ và ràng buộc trạng thái",
            "phongkham.dao: truy cập CSDL và truy vấn có tham số",
            "phongkham.DTO: mô hình dữ liệu trao đổi giữa các tầng",
            "phongkham.db: cấu hình pool kết nối",
            "phongkham.Utils: session, xuất PDF/CSV, helper dùng chung",
        ],
        "Lợi ích",
        [
            "Dễ mở rộng từng module độc lập",
            "Giảm phụ thuộc chéo giữa GUI và SQL",
            "Dễ kiểm tra lỗi theo tầng",
            "Thuận lợi chia việc theo nhóm",
            "Giữ được tính nhất quán dữ liệu toàn hệ thống",
        ],
    )

    # 7 - RBAC
    add_slide(
        prs,
        "RBAC động: đăng nhập là nạp quyền",
        bullets=[
            "Flow: Login -> PermissionBUS.loadPermission(userId) -> Session lưu quyền dạng UPPERCASE.",
            "MainFrame.canAccess(route) quyết định cho vào route hay không.",
            "Sidebar.buildMenu() dựng menu theo quyền hiện có của người dùng.",
            "Admin đổi quyền trong màn hình Phân quyền: tác dụng ngay sau đăng nhập lại, không sửa code GUI.",
            "Route đặc biệt QL_DUYET_LICH_LAM dùng mode ADMIN để tăng độ kiểm soát.",
        ],
    )

    # 8 - Permission map
    slide = add_slide(prs, "Bản đồ quyền trong hệ thống", "Ví dụ quyền thực tế đang dùng")
    add_two_columns(
        slide,
        "Doctor + Pharmacy",
        [
            "LICHLAMVIEC_VIEW / LICHLAMVIEC_MANAGE",
            "LICHKHAM_VIEW / LICHKHAM_MANAGE",
            "HOSO_VIEW / HOSO_MANAGE",
            "HOADONKHAM_VIEW / HOADONKHAM_MANAGE",
            "THUOC_VIEW / THUOC_MANAGE",
            "PHIEUNHAP_VIEW / PHIEUNHAP_MANAGE",
            "HOADONTHUOC_VIEW / HOADONTHUOC_CREATE / HOADONTHUOC_MANAGE",
        ],
        "Admin",
        [
            "DASHBOARD_VIEW",
            "USER_VIEW / USER_MANAGE",
            "BACSI_VIEW / BACSI_MANAGE",
            "KHOA_VIEW / KHOA_MANAGE",
            "GOIDICHVU_VIEW / GOIDICHVU_MANAGE",
            "ROLE_PERMISSION_VIEW / ROLE_PERMISSION_MANAGE",
            "Guest chỉ có 2 route công khai: DAT_LICH, MUA_THUOC",
        ],
    )

    # 9 - Workflow guest
    add_slide(
        prs,
        "Luồng 1: Khách hàng đặt lịch khám",
        bullets=[
            "B1. Chọn gói dịch vụ -> hệ thống lọc bác sĩ theo khoa tương ứng.",
            "B2. Chọn ngày/ca chỉ từ lịch làm việc đã duyệt.",
            "B3. Tạo lịch khám trạng thái CHO_XAC_NHAN.",
            "B4. Có thể tra cứu hồ sơ theo CCCD và in phiếu đặt lịch PDF.",
            "Giá trị mang lại: giảm xung đột lịch và chuẩn hóa đầu vào cho bác sĩ.",
        ],
    )

    # 10 - Workflow doctor
    add_slide(
        prs,
        "Luồng 2: Bác sĩ xử lý khám bệnh",
        bullets=[
            "B1. Bác sĩ xem lịch khám được gán và xác nhận/hủy theo quy tắc trạng thái.",
            "B2. Sau khám, lập hồ sơ bệnh án: triệu chứng, chẩn đoán, kết luận, lời dặn.",
            "B3. Kê đơn thuốc: kiểm tra số lượng tồn trước khi lưu.",
            "B4. Hóa đơn khám được quản lý trạng thái thanh toán rõ ràng.",
            "Ràng buộc quan trọng: không cho sửa/xóa sai trạng thái đã hoàn tất.",
        ],
    )

    # 11 - Workflow pharmacy
    add_slide(
        prs,
        "Luồng 3: Nhà thuốc và kho",
        bullets=[
            "B1. Tạo phiếu nhập: chọn NCC, thêm thuốc theo số lô và hạn sử dụng.",
            "B2. Xác nhận nhập kho làm tăng tồn và tạo dữ liệu truy vết lô.",
            "B3. Tạo hóa đơn thuốc từ đơn hoặc bán lẻ cho khách không đơn.",
            "B4. Quy trình thanh toán/lấy thuốc theo trạng thái; hỗ trợ xuất CSV truy vết.",
            "Điểm mạnh: mô hình lô thuốc hỗ trợ kiểm soát hạn dùng và tồn kho thực tế.",
        ],
    )

    # 12 - Workflow admin
    add_slide(
        prs,
        "Luồng 4: Quản trị và điều hành",
        bullets=[
            "Dashboard: doanh thu 6 tháng, lịch khám hôm nay, cảnh báo thuốc cận hạn/tồn thấp.",
            "Quản lý tài khoản: thêm/sửa/khóa-mở/reset mật khẩu.",
            "Quản lý bác sĩ/khoa/gói dịch vụ và duyệt lịch làm.",
            "Phân quyền role-permission trực tiếp trên dữ liệu để thay đổi quyền vận hành tức thì.",
        ],
    )

    # 13 - Data model
    slide = add_slide(prs, "Thiết kế dữ liệu logic", "Nhóm bảng chính và ý nghĩa")
    add_two_columns(
        slide,
        "Nhóm bảng",
        [
            "Users, Roles, Permissions, RolePermissions",
            "BacSi, Khoa, LichLamViec, LichKham, HoSoBenhAn",
            "DonThuoc, CTDonThuoc, HoaDonKham",
            "Thuoc, NhaCungCap, PhieuNhap, CTPhieuNhap",
            "HoaDonThuoc, CTHDThuoc",
        ],
        "Ý nghĩa",
        [
            "Phân quyền linh hoạt theo RBAC",
            "Chuỗi nghiệp vụ khám bệnh đầy đủ từ lịch -> bệnh án",
            "Kê đơn và thanh toán dịch vụ khám",
            "Nhập kho theo lô-hạn dùng, giữ lịch sử biến động",
            "Bán thuốc và truy vết xuất kho",
        ],
    )

    # 14 - Data constraints
    add_slide(
        prs,
        "Ràng buộc dữ liệu tiêu biểu",
        bullets=[
            "UNIQUE(MaBacSi, NgayLam, CaLam): ngăn đăng ký trùng ca làm.",
            "CHECK(CaLam in Sang/Chieu/Toi): chuẩn hóa giá trị lịch làm.",
            "UNIQUE(MaPhieuNhap, MaThuoc, SoLo, HanSuDung): không trùng lô nhập.",
            "SoLuongTon >= 0 và kiểm tra nghiệp vụ trước thao tác trừ tồn.",
            "Hóa đơn đã thanh toán không cho xóa để bảo toàn tính toàn vẹn nghiệp vụ.",
        ],
    )

    # 15 - Technical highlights
    slide = add_slide(prs, "Điểm nhấn kỹ thuật")
    add_two_columns(
        slide,
        "Nền tảng kỹ thuật",
        [
            "Java Swing + FlatLaf cho desktop UI hiện đại",
            "CardLayout cho điều hướng panel theo route",
            "HikariCP pooling (max 12, min idle 2)",
            "PreparedStatement và cache câu lệnh",
            "Xuất PDF/CSV phục vụ vận hành thực tế",
        ],
        "Giá trị thực tiễn",
        [
            "Tốc độ phản hồi ổn định khi thao tác nhiều màn hình",
            "Giảm lỗi SQL injection và lỗi nhập liệu",
            "Dễ bảo trì vì tách tầng rõ ràng",
            "Bám sát nghiệp vụ thực tế phòng khám + nhà thuốc",
            "Có thể nâng cấp dần lên web/service trong tương lai",
        ],
    )

    # 16 - Strength and limitations
    slide = add_slide(prs, "Đánh giá khách quan", "Điểm mạnh và điểm cần cải tiến")
    add_two_columns(
        slide,
        "Điểm mạnh",
        [
            "Phân lớp rõ, module hóa tốt, code dễ đọc",
            "RBAC động đổi quyền không cần build lại",
            "Quản lý theo trạng thái nghiệp vụ chặt chẽ",
            "Truy vết lô thuốc hỗ trợ kiểm toán",
            "Có sẵn nhiều báo cáo xuất file",
        ],
        "Cần cải tiến",
        [
            "Thiếu transaction bao bọc các thao tác nhiều bước",
            "Chưa có audit log chi tiết cho thay đổi quyền/dữ liệu",
            "Error logging chưa tập trung theo chuẩn production",
            "UI và business logic còn dính nhau ở vài màn hình",
            "Cần bổ sung test tự động sâu hơn cho BUS/DAO",
        ],
    )

    # 17 - Demo plan
    add_slide(
        prs,
        "Kịch bản demo đề xuất (15-20 phút)",
        bullets=[
            "Phần 1 (3 phút): Đăng nhập 4 vai trò để thấy menu thay đổi theo quyền.",
            "Phần 2 (5 phút): Guest đặt lịch -> Doctor xác nhận -> lập bệnh án + đơn.",
            "Phần 3 (5 phút): Pharmacist nhập kho theo lô -> bán thuốc -> xuất CSV truy vết.",
            "Phần 4 (3 phút): Admin đổi quyền ngay trên Phân quyền và kiểm tra hiệu lực.",
            "Phần 5 (2 phút): Mở Dashboard, giải thích số liệu vận hành và cảnh báo.",
        ],
    )

    # 18 - Q&A defense
    add_slide(
        prs,
        "Câu hỏi phản biện thường gặp",
        bullets=[
            "Vì sao chọn desktop Swing thay vì web ngay từ đầu?",
            (1, "Trả lời: phù hợp đồ án quản trị nội bộ, triển khai nhanh, dễ kiểm soát local dataflow."),
            "RBAC động có đảm bảo an toàn truy cập route?",
            (1, "Trả lời: có chặn 2 lớp: ẩn menu ở Sidebar và chặn route tại MainFrame.canAccess."),
            "Làm sao đảm bảo dữ liệu tồn kho đúng?",
            (1, "Trả lời: ràng buộc SQL + kiểm tra BUS trước ghi + truy vết theo lô/hạn dùng."),
            "Hướng nâng cấp tương lai?",
            (1, "Trả lời: tách BUS/DAO thành service layer, bổ sung API và unit/integration tests."),
        ],
    )

    # 19 - Team contribution
    add_slide(
        prs,
        "Tổ chức triển khai và phân công",
        bullets=[
            "Thiết kế dữ liệu và migration script: chuẩn hóa quan hệ bảng và constraint.",
            "Phát triển nghiệp vụ theo module: Guest, Doctor, Pharmacy, Admin.",
            "Tích hợp tiện ích xuất báo cáo: PDF/CSV cho vận hành.",
            "Kiểm thử và rà soát logic trạng thái theo từng luồng.",
            "Tài liệu hóa: báo cáo phân tích, ca kiểm thử và slide thuyết trình.",
        ],
    )

    # 20 - Closing
    slide = add_slide(
        prs,
        "Kết luận",
        "Hệ thống đáp ứng tốt yêu cầu quản lý phòng khám đa vai trò với khả năng mở rộng cao",
        dark=True,
        bullets=[
            "Đã giải quyết bài toán lõi: khám bệnh, kê đơn, bán thuốc, phân quyền, báo cáo.",
            "Kiến trúc rõ ràng giúp hệ thống dễ bảo trì và nâng cấp.",
            "RBAC động là điểm nhấn tạo khác biệt khi triển khai thực tế.",
            "Sẵn sàng cho giai đoạn cải tiến: transaction, audit log, test automation, API hóa.",
            "Xin cảm ơn thầy/cô và các bạn đã lắng nghe.",
        ],
    )
    thanks = slide.shapes.add_textbox(Inches(8.0), Inches(5.85), Inches(4.8), Inches(0.7))
    ttf = thanks.text_frame
    ttf.clear()
    tp = ttf.paragraphs[0]
    tp.text = "Q&A"
    tp.font.name = "Calibri"
    tp.font.size = Pt(40)
    tp.font.bold = True
    tp.font.color.rgb = ACCENT_2
    tp.alignment = PP_ALIGN.RIGHT

    prs.save(OUTPUT_FILE)


if __name__ == "__main__":
    build_presentation()
    print(f"Generated: {OUTPUT_FILE}")
